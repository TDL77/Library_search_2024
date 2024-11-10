import streamlit as st
import os
from pathlib import Path
import pypdf
import pandas as pd
import json
import openpyxl
from pptx import Presentation
import base64
from io import BytesIO
from datetime import datetime, timedelta
import tkinter as tk
from tkinter import filedialog
import subprocess
import platform

from main import Config, RAGSystem

config = Config()
rag_system = RAGSystem(config)

# Supported file formats
SUPPORTED_FORMATS = {'.docx', '.xlsx', '.pptx', '.pdf', '.csv', '.md', '.json'}

def get_all_files(directory):
    """Recursively get all files from directory and its subdirectories"""
    files_list = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            rel_path = os.path.relpath(file_path, directory)
            if os.path.splitext(file)[1].lower() in SUPPORTED_FORMATS:
                files_list.append((rel_path, file_path))
    return files_list

# Инициализация session state
if "messages" not in st.session_state:
    st.session_state.messages = []

if "selected_folder" not in st.session_state:
    st.session_state.selected_folder = "./data/"

# Настройка конфигурации страницы
st.set_page_config(
    page_title="Document Assistant",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Стили CSS
st.markdown("""
    <style>
    .main {
        padding: 2rem;
    }
    .stButton>button {
        width: 100%;
        margin-top: 1rem;
        background-color: #4CAF50;
        color: white;
    }
    
    .chat-message {
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
        max-width: 80%;
        color: white;
    }
    
    .user-message {
        background-color: #2C3E50;
        margin-left: auto;
        margin-right: 0;
    }
    
    .bot-message {
        background-color: #34495E;
        margin-left: 0;
        margin-right: auto;
    }
    
    .document-viewer {
        height: 50vh;
        overflow-y: auto;
        border: 1px solid #ddd;
        padding: 1rem;
        margin-top: 1rem;
        background-color: #2C3E50;
        color: white;
        border-radius: 0.5rem;
    }
    
    .document-title {
        color: white;
        background-color: #2C3E50;
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    
    .dataframe {
        color: white !important;
        background-color: #34495E !important;
    }
    
    .dataframe th {
        background-color: #2C3E50 !important;
        color: white !important;
    }
    
    .dataframe td {
        background-color: #34495E !important;
        color: white !important;
    }
    
    .json-content {
        background-color: #2C3E50 !important;
        color: white !important;
        border-radius: 0.5rem;
        padding: 1rem;
    }
    
    iframe {
        background-color: white;
        border-radius: 0.5rem;
        margin-top: 1rem;
    }

    .stats-container {
        background-color: #2E8B57;
        color: white;
        border-radius: 10px;
        padding: 15px;
        margin: 10px 0;
    }
    
    .filter-section {
        margin-top: 50px;
        padding: 20px;
        background-color: transparent;
        border-radius: 10px;
    }
    
    .filter-header {
        color: white;
        font-size: 24px;
        margin-bottom: 20px;
        padding-bottom: 10px;
        border-bottom: 2px solid #2E8B57;
    }
    </style>
    """, unsafe_allow_html=True)

def read_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        pdf_reader = pypdf.PdfReader(file)
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
    return text

def read_xlsx(file_path):
    df = pd.read_excel(file_path)
    return df.to_html()

def read_csv(file_path):
    df = pd.read_csv(file_path)
    return df.to_html()

def read_md(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        return file.read()

def read_json(file_path):
    with open(file_path, 'r', encoding='utf-8') as file:
        data = json.load(file)
        return json.dumps(data, indent=2)

def read_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def display_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension == '.pdf':
            with open(file_path, "rb") as f:
                base64_pdf = base64.b64encode(f.read()).decode('utf-8')
            pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="600" type="application/pdf"></iframe>'
            st.markdown(pdf_display, unsafe_allow_html=True)
        
        elif file_extension in ['.xlsx', '.csv']:
            if file_extension == '.xlsx':
                df = pd.read_excel(file_path)
            else:
                df = pd.read_csv(file_path)
            st.markdown('<div class="document-content">', unsafe_allow_html=True)
            st.dataframe(df, height=400)
            st.markdown('</div>', unsafe_allow_html=True)
        
        elif file_extension == '.json':
            content = read_json(file_path)
            st.markdown('<div class="json-content">', unsafe_allow_html=True)
            st.json(json.loads(content))
            st.markdown('</div>', unsafe_allow_html=True)
        
        elif file_extension == '.md':
            content = read_md(file_path)
            st.markdown(f'<div class="document-content">{content}</div>', unsafe_allow_html=True)
        
        elif file_extension == '.pptx':
            content = read_pptx(file_path)
            st.markdown(f'<div class="document-content">{content}</div>', unsafe_allow_html=True)
        
    except Exception as e:
        st.error(f"Ошибка при чтении файла: {str(e)}")

# Боковая панель
with st.sidebar:
    st.title("📚 Настройки")
    
    # Set default folder path
    folder_path = st.text_input("Введите путь к папке с документами:", 
                               value="./data/",
                               help="Укажите полный путь к папке с документами")
    
    if folder_path and os.path.exists(folder_path):
        st.session_state.selected_folder = folder_path
        st.success("✅ Папка успешно выбрана!")
        
        # Get all files recursively
        all_files = get_all_files(folder_path)
        
        # Show folder statistics
        st.markdown("### 📊 Анализ содержимого папки")
        stats_folder_name = Path(folder_path).name
        last_modified = datetime.fromtimestamp(os.path.getmtime(folder_path)).strftime('%Y-%m-%d %H:%M:%S')
        total_files = len(all_files)
        
        st.markdown(
            f"""<div class='stats-container'>
            <strong>📁 Selected:</strong> {stats_folder_name}<br/>
            <strong>📄 Total Files:</strong> {total_files}<br/>
            <strong>🕒 Last Modified:</strong> {last_modified}
            </div>""", 
            unsafe_allow_html=True
        )
        
        # Document filters
        st.markdown("### 🔍 Фильтр документов")
        st.markdown("<div class='filter-container'>", unsafe_allow_html=True)
        
        # File type selection
        available_extensions = set()
        for _, file_path in all_files:
            ext = os.path.splitext(file_path)[1].lower()
            if ext in SUPPORTED_FORMATS:
                available_extensions.add(ext)
                
        selected_extensions = st.multiselect(
            "Выберите тип файлов:",
            options=sorted(available_extensions),
            default=None
        )
        
        # Date filter
        date_filter = st.selectbox(
            "Выберите актуальность данных:",
            options=["За все время", "Последний месяц", "За последние 3 месяца", 
                    "За последние 6 месяцев", "За последний год", "Более чем год"],
            index=0
        )
        
        if st.button("Применить фильтры"):
            if selected_extensions:
                filtered_files = [(rel_path, file_path) for rel_path, file_path in all_files 
                                if os.path.splitext(file_path)[1].lower() in selected_extensions]
                
                # Apply date filter
                if date_filter != "За все время":
                    now = datetime.now()
                    date_filters = {
                        "Последний месяц": timedelta(days=30),
                        "За последние 3 месяца": timedelta(days=90),
                        "За последние 6 месяцев": timedelta(days=180),
                        "За последний год": timedelta(days=365),
                        "Более чем год": None
                    }
                    
                    filtered_files = [
                        (rel_path, file_path) for rel_path, file_path in filtered_files
                        if date_filters[date_filter] is None or 
                        (now - datetime.fromtimestamp(os.path.getctime(file_path))) 
                        <= date_filters[date_filter]
                    ]
                
                st.success(f"Найдено {len(filtered_files)} документов")
        
        # Show file types found
        st.markdown("#### 📌 Найдены типы файлов:")
        extension_counts = {}
        for _, file_path in all_files:
            ext = os.path.splitext(file_path)[1].lower()
            if ext in SUPPORTED_FORMATS:
                extension_counts[ext] = extension_counts.get(ext, 0) + 1
        
        for ext, count in sorted(extension_counts.items()):
            st.markdown(f"<div style='color: #2E8B57;'>• {ext}: {count} files</div>", 
                      unsafe_allow_html=True)

                  # Загрузка файлов
        uploaded_files = st.file_uploader("Загрузите файлы", type=list(SUPPORTED_FORMATS), accept_multiple_files=True)
        if uploaded_files:
            for uploaded_file in uploaded_files:
                save_path = os.path.join(st.session_state.selected_folder, uploaded_file.name)
                with open(save_path, "wb") as f:
                    f.write(uploaded_file.getbuffer())
            st.success("Файлы успешно загружены!")

# Main interface
col1, col2 = st.columns([1, 1])

with col1:
    st.title("💬 Чат с ассистентом")
    
    user_input = st.text_area("Ваш вопрос:", height=100)
    
    if st.button("Отправить вопрос"):
        if user_input:
            st.session_state.messages.append({"role": "user", "content": user_input})
            # Get context and generate response using RAG system
            context, docs, confidence, metadata = rag_system.get_relevant_content(user_input)
            response = rag_system.generate_response(user_input, context)
            st.session_state.messages.append({"role": "assistant", "content": response})
    
    # Display message history
    st.subheader("История диалога:")
    for message in st.session_state.messages:
        if message["role"] == "user":
            st.markdown(f"""<div class="chat-message user-message">
                          <b>Вы:</b><br>{message["content"]}</div>""", 
                       unsafe_allow_html=True)
        else:
            st.markdown(f"""<div class="chat-message bot-message">
                          <b>Ассистент:</b><br>{message["content"]}</div>""", 
                       unsafe_allow_html=True)

# Document viewer
with col2:
    st.markdown('<h2 class="document-title">📄 Просмотр документа</h2>', unsafe_allow_html=True)
    
    if st.session_state.selected_folder:
        all_files = get_all_files(st.session_state.selected_folder)
        
        if all_files:
            file_options = [rel_path for rel_path, _ in all_files]
            
            # Create a default option for when no document is selected
            default_text = "Путь до релевантного файла..."
            
            # If metadata exists and has a source, use it as the selected value
            selected_value = default_text
            if 'metadata' in locals() and metadata and 'source' in metadata:
                selected_value = metadata['source']
            
            # Display the selectbox with either default text or metadata source
            st.selectbox("Путь до релевантного файла...", 
                        [selected_value],
                        disabled=True)
            
            # Display the content from docs[0].page_content
            st.markdown('<div class="document-container">', unsafe_allow_html=True)
            st.markdown('<h3 class="document-title">Релевантное содержимое:</h3>', 
                      unsafe_allow_html=True)
            
            # Create a custom container for the content with better formatting
            st.markdown("""
                <style>
                .content-display {
                    background-color: #2C3E50;
                    color: white;
                    padding: 20px;
                    border-radius: 10px;
                    margin-top: 10px;
                    white-space: pre-wrap;
                    font-family: monospace;
                    line-height: 1.5;
                    max-height: 600px;
                    overflow-y: auto;
                }
                </style>
            """, unsafe_allow_html=True)
            
            # Display the content with line breaks preserved
            if 'docs' in locals() and docs and len(docs) > 0:
                content = docs[0].page_content
                formatted_content = content.replace('\n', '<br>')
                st.markdown(f'<div class="content-display">{formatted_content}</div>', 
                          unsafe_allow_html=True)
            else:
                st.info("Задайте вопрос, чтобы увидеть релевантное содержимое документов")
            
            st.markdown('</div>', unsafe_allow_html=True)
        else:
            st.warning("В выбранной папке нет поддерживаемых документов")
    else:
        st.info("👈 Сначала выберите папку с документами в боковой панели")